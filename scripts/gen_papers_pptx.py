#!/home/hermes/.venv/bin/python3
"""Generate a .pptx deck from /tmp/papers.json. Output: /tmp/ai-papers-gist.pptx"""
import json, sys
from pathlib import Path
from datetime import date

INPUT = "/tmp/papers.json"
OUTPUT = "/tmp/ai-papers-gist.pptx"

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
from pptx.util import Inches, Pt

papers = json.loads(Path(INPUT).read_text())

def _to_bullets(p):
    if "bullets" in p and isinstance(p["bullets"], list):
        return p["bullets"]
    parts = []
    if p.get("authors"):
        parts.append(f"Authors: {p['authors']}")
    words, line = (p.get("summary","")).split(), []
    for w in words:
        line.append(w)
        if len(" ".join(line)) > 110:
            parts.append(" ".join(line)); line = []
    if line: parts.append(" ".join(line))
    return parts[:4] or ["No summary available."]

prs = Presentation()
title_layout = prs.slide_layouts[0]
bullet_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = "AI Papers — Gist Summary"
slide.placeholders[1].text = str(date.today())

for p in papers:
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = p["title"]
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for j, bullet in enumerate(_to_bullets(p)):
        if j == 0: tf.paragraphs[0].text = bullet
        else:
            para = tf.add_paragraph(); para.text = bullet

prs.save(OUTPUT)
print(f"PPTX written to {OUTPUT} ({len(papers)} papers)")
